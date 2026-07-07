from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from parser import DesignDNA, PresentationData, Slide

# 16:9 widescreen slide dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Shape type constant for a rectangle (MSO_AUTO_SHAPE_TYPE.RECTANGLE = 1)
_RECT = 1


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_bg(slide, color_hex: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def _rect(slide, left, top, width, height, color_hex: str):
    """Add a solid rectangle with no border."""
    shape = slide.shapes.add_shape(
        _RECT, int(left), int(top), int(width), int(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color_hex)
    shape.line.fill.background()
    return shape


def _text(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    size: float,
    color_hex: str,
    bold: bool = False,
    italic: bool = False,
    align: str = "left",
    font: str = "Calibri",
):
    """Add a text box with a single paragraph."""
    txBox = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.color.rgb = hex_to_rgb(color_hex)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    return txBox


# ---------------------------------------------------------------------------
# Layout renderers
# ---------------------------------------------------------------------------

def _cover(slide, s: Slide, dna: DesignDNA):
    W, H = SLIDE_W, SLIDE_H
    _set_bg(slide, dna.bg_dark)

    # Accent strip top + bottom
    _rect(slide, 0, 0, W, H * 0.012, dna.accent)
    _rect(slide, 0, H * 0.988, W, H * 0.012, dna.accent)

    # Left accent vertical bar
    _rect(slide, 0, 0, W * 0.007, H, dna.accent)

    # Decorative side panel (muted accent — mix accent with bg_dark)
    panel_color = _blend(dna.accent, dna.bg_dark, 0.18)
    _rect(slide, W * 0.62, H * 0.1, W * 0.38, H * 0.8, panel_color)

    # Title
    _text(
        slide, s.title,
        W * 0.08, H * 0.22, W * 0.78, H * 0.36,
        size=46, color_hex=dna.text_dark, bold=True, align="left", font=dna.font,
    )
    # Accent separator line
    _rect(slide, W * 0.08, H * 0.61, W * 0.4, H * 0.007, dna.accent)

    # Subtitle
    if s.subtitle:
        _text(
            slide, s.subtitle,
            W * 0.08, H * 0.65, W * 0.78, H * 0.2,
            size=22, color_hex=dna.text_dark, align="left", font=dna.font,
        )


def _section(slide, s: Slide, dna: DesignDNA):
    W, H = SLIDE_W, SLIDE_H
    _set_bg(slide, dna.accent)

    # Dark side panel (right 40%)
    _rect(slide, W * 0.6, 0, W * 0.4, H, _blend(dna.bg_dark, dna.accent, 0.35))

    # Top and bottom thin dark strips for visual rhythm with other slides
    _rect(slide, 0, 0, W, H * 0.012, dna.bg_dark)
    _rect(slide, 0, H * 0.988, W, H * 0.012, dna.bg_dark)

    # Title
    _text(
        slide, s.title,
        W * 0.07, H * 0.28, W * 0.55, H * 0.44,
        size=44, color_hex="#FFFFFF", bold=True, align="left", font=dna.font,
    )

    if s.subtitle:
        _text(
            slide, s.subtitle,
            W * 0.07, H * 0.73, W * 0.55, H * 0.15,
            size=20, color_hex="#FFFFFF", align="left", font=dna.font,
        )


def _bullets(slide, s: Slide, dna: DesignDNA):
    W, H = SLIDE_W, SLIDE_H
    _set_bg(slide, dna.bg_light)

    # Left accent bar
    _rect(slide, 0, 0, W * 0.007, H, dna.accent)

    # Title strip (dark background)
    _rect(slide, 0, 0, W, H * 0.22, dna.bg_dark)
    _text(
        slide, s.title,
        W * 0.04, H * 0.04, W * 0.92, H * 0.16,
        size=30, color_hex=dna.text_dark, bold=True, align="left", font=dna.font,
    )
    # Accent underline
    _rect(slide, 0, H * 0.22, W, H * 0.006, dna.accent)

    # Bullet items
    y = H * 0.3
    step = H * 0.13
    for bullet in (s.bullets or [])[:5]:
        _text(
            slide, f"•  {bullet}",
            W * 0.05, y, W * 0.9, step * 0.95,
            size=20, color_hex=dna.text_light, align="left", font=dna.font,
        )
        y += step


def _quote(slide, s: Slide, dna: DesignDNA):
    W, H = SLIDE_W, SLIDE_H
    _set_bg(slide, dna.bg_dark)

    # Accent top and bottom strips
    _rect(slide, 0, 0, W, H * 0.012, dna.accent)
    _rect(slide, 0, H * 0.988, W, H * 0.012, dna.accent)

    # Large decorative quotation mark (visually prominent)
    _text(
        slide, "“",
        W * 0.03, H * 0.01, W * 0.2, H * 0.35,
        size=100, color_hex=dna.accent, bold=True, align="left", font=dna.font,
    )

    # Quote body
    quote_body = s.quote or s.title
    _text(
        slide, quote_body,
        W * 0.1, H * 0.28, W * 0.8, H * 0.42,
        size=26, color_hex=dna.text_dark, italic=True, align="center", font=dna.font,
    )

    # Accent separator
    _rect(slide, W * 0.35, H * 0.73, W * 0.3, H * 0.006, dna.accent)

    # Attribution (title used as source when quote is the body)
    source = s.title if s.quote else ""
    if source:
        _text(
            slide, f"— {source}",
            W * 0.1, H * 0.77, W * 0.8, H * 0.1,
            size=18, color_hex=dna.text_dark, align="center", font=dna.font,
        )


def _stat(slide, s: Slide, dna: DesignDNA):
    W, H = SLIDE_W, SLIDE_H
    _set_bg(slide, dna.bg_light)

    # Left half in accent color
    _rect(slide, 0, 0, W * 0.46, H, dna.accent)

    # Title strip spans full width
    _rect(slide, 0, 0, W, H * 0.18, dna.bg_dark)
    _text(
        slide, s.title,
        W * 0.04, H * 0.03, W * 0.92, H * 0.14,
        size=28, color_hex=dna.text_dark, bold=True, align="center", font=dna.font,
    )

    # Thin accent divider between left/right halves (below title strip)
    _rect(slide, W * 0.455, H * 0.18, W * 0.008, H * 0.82, _blend(dna.bg_dark, dna.accent, 0.4))

    # Large stat number (left half)
    _text(
        slide, s.stat_value or "—",
        W * 0.0, H * 0.2, W * 0.46, H * 0.65,
        size=72, color_hex="#FFFFFF", bold=True, align="center", font=dna.font,
    )

    # Stat label (right half)
    _text(
        slide, s.stat_label or "",
        W * 0.52, H * 0.3, W * 0.44, H * 0.5,
        size=22, color_hex=dna.text_light, align="left", font=dna.font,
    )


def _two_col(slide, s: Slide, dna: DesignDNA):
    W, H = SLIDE_W, SLIDE_H
    _set_bg(slide, dna.bg_light)

    # Title strip
    _rect(slide, 0, 0, W, H * 0.22, dna.bg_dark)
    _text(
        slide, s.title,
        W * 0.04, H * 0.04, W * 0.92, H * 0.15,
        size=30, color_hex=dna.text_dark, bold=True, align="center", font=dna.font,
    )
    _rect(slide, 0, H * 0.22, W, H * 0.006, dna.accent)

    # Center divider
    _rect(slide, W * 0.492, H * 0.25, W * 0.008, H * 0.68, dna.accent)

    left_d = s.left or {}
    right_d = s.right or {}

    # Left column
    if left_d.get("title"):
        _text(
            slide, left_d["title"],
            W * 0.04, H * 0.27, W * 0.43, H * 0.12,
            size=22, color_hex=dna.accent, bold=True, align="left", font=dna.font,
        )
    if left_d.get("body"):
        _text(
            slide, left_d["body"],
            W * 0.04, H * 0.41, W * 0.43, H * 0.52,
            size=18, color_hex=dna.text_light, align="left", font=dna.font,
        )

    # Right column
    if right_d.get("title"):
        _text(
            slide, right_d["title"],
            W * 0.52, H * 0.27, W * 0.44, H * 0.12,
            size=22, color_hex=dna.accent, bold=True, align="left", font=dna.font,
        )
    if right_d.get("body"):
        _text(
            slide, right_d["body"],
            W * 0.52, H * 0.41, W * 0.44, H * 0.52,
            size=18, color_hex=dna.text_light, align="left", font=dna.font,
        )


def _closing(slide, s: Slide, dna: DesignDNA):
    W, H = SLIDE_W, SLIDE_H
    _set_bg(slide, dna.bg_dark)

    # Accent strips top + bottom (mirrors cover)
    _rect(slide, 0, 0, W, H * 0.012, dna.accent)
    _rect(slide, 0, H * 0.988, W, H * 0.012, dna.accent)

    # Right vertical bar (mirror of cover's left bar)
    _rect(slide, W * 0.993, 0, W * 0.007, H, dna.accent)

    # Decorative left panel (mirrors cover's right panel)
    panel_color = _blend(dna.accent, dna.bg_dark, 0.18)
    _rect(slide, 0, H * 0.1, W * 0.38, H * 0.8, panel_color)

    # Title
    _text(
        slide, s.title,
        W * 0.08, H * 0.22, W * 0.84, H * 0.36,
        size=46, color_hex=dna.text_dark, bold=True, align="center", font=dna.font,
    )
    # Accent separator
    _rect(slide, W * 0.3, H * 0.61, W * 0.4, H * 0.007, dna.accent)

    # CTA subtitle
    if s.subtitle:
        _text(
            slide, s.subtitle,
            W * 0.08, H * 0.66, W * 0.84, H * 0.2,
            size=22, color_hex=dna.accent, align="center", font=dna.font,
        )


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _blend(hex_a: str, hex_b: str, ratio_a: float) -> str:
    """Blend two hex colors: ratio_a fraction of hex_a, rest of hex_b."""
    a = hex_a.lstrip("#")
    b = hex_b.lstrip("#")
    r = int(int(a[0:2], 16) * ratio_a + int(b[0:2], 16) * (1 - ratio_a))
    g = int(int(a[2:4], 16) * ratio_a + int(b[2:4], 16) * (1 - ratio_a))
    bv = int(int(a[4:6], 16) * ratio_a + int(b[4:6], 16) * (1 - ratio_a))
    return f"#{r:02X}{g:02X}{bv:02X}"


_RENDERERS = {
    "cover": _cover,
    "section": _section,
    "bullets": _bullets,
    "quote": _quote,
    "stat": _stat,
    "two_col": _two_col,
    "closing": _closing,
}


def render_pptx(data: PresentationData, output_path: str):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Index 6 is the blank layout in a default python-pptx Presentation
    blank = prs.slide_layouts[6]

    for slide_data in data.slides:
        slide = prs.slides.add_slide(blank)
        renderer = _RENDERERS.get(slide_data.layout, _bullets)
        renderer(slide, slide_data, data.design_dna)

    prs.save(output_path)
